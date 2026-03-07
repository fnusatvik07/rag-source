"""
Structured slide-by-slide PPTX extraction for RAG.

Demonstrates how to:
  - Parse each slide into a structured dict (slide_number, title, body_text,
    table_data, notes)
  - Convert structured slides into RAG-ready text chunks (one per slide)
  - Build slide summaries that preserve context
  - Apply sentence-based chunking on the full extracted text

Run:
    uv run python unstructured_documents/03_pptx/02_slide_structured_extraction.py
"""

import json
import sys
from pathlib import Path

# --- shared chunking import ------------------------------------------------
sys.path.insert(0, str(Path(__file__).resolve().parents[2]))
from pptx import Presentation

from unstructured_documents.shared.chunking import (
    chunk_by_sentences,
    preview_chunks,
)

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
SAMPLE_DIR = Path(__file__).resolve().parent / "sample_docs"
PPTX_PATH = SAMPLE_DIR / "presentation.pptx"


# ---------------------------------------------------------------------------
# Structured extraction
# ---------------------------------------------------------------------------


def _collect_body_text(shape) -> list[str]:
    """
    Recursively collect body text from a shape (skipping titles).

    Handles plain text frames and group shapes.
    """
    texts = []

    # Group shape — recurse
    if shape.shape_type is not None and shape.shape_type == 6:
        for child in shape.shapes:
            texts.extend(_collect_body_text(child))
        return texts

    # Text frame (but not the slide title placeholder)
    if shape.has_text_frame:
        # placeholder idx 0 is the title — skip it here
        try:
            if shape.placeholder_format is not None and shape.placeholder_format.idx == 0:
                return texts
        except ValueError:
            pass  # shape is not a placeholder — that's fine, extract its text
        text = shape.text_frame.text.strip()
        if text:
            texts.append(text)

    return texts


def _collect_table_data(shape) -> list[list[list[str]]] | None:
    """Return table rows if the shape is a table, else None."""
    if not shape.has_table:
        return None
    rows = []
    for row in shape.table.rows:
        rows.append([cell.text.strip() for cell in row.cells])
    return rows


def extract_structured_slides(pptx_path: Path) -> list[dict]:
    """
    Parse a PPTX into a list of structured slide dictionaries.

    Each dict contains:
      - slide_number (int)
      - title (str)
      - body_text (str)       — all non-title text joined
      - table_data (list)     — list of tables, each a 2-D list
      - notes (str)           — speaker notes
    """
    prs = Presentation(str(pptx_path))
    structured = []

    for idx, slide in enumerate(prs.slides, start=1):
        # --- Title ---
        title = ""
        if slide.shapes.title is not None:
            title = slide.shapes.title.text.strip()

        # --- Body text and tables ---
        body_parts = []
        tables = []

        for shape in slide.shapes:
            # Tables
            tbl = _collect_table_data(shape)
            if tbl is not None:
                tables.append(tbl)
                continue

            # Body text (non-title)
            body_parts.extend(_collect_body_text(shape))

        body_text = "\n".join(body_parts)

        # --- Speaker notes ---
        notes = ""
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()

        structured.append(
            {
                "slide_number": idx,
                "title": title,
                "body_text": body_text,
                "table_data": tables,
                "notes": notes,
            }
        )

    return structured


# ---------------------------------------------------------------------------
# RAG-ready conversion
# ---------------------------------------------------------------------------


def table_to_text(table: list[list[str]]) -> str:
    """
    Convert a 2-D table into a readable text block.

    Uses the first row as headers and formats remaining rows as
    "header: value" pairs for better retrieval relevance.
    """
    if not table or len(table) < 2:
        return ""

    headers = table[0]
    lines = []
    for row in table[1:]:
        parts = []
        for h, v in zip(headers, row):
            parts.append(f"{h}: {v}")
        lines.append("; ".join(parts))
    return "\n".join(lines)


def slides_to_rag_chunks(slides: list[dict], include_notes: bool = True) -> list[dict]:
    """
    Convert structured slides into RAG-ready text chunks.

    Each chunk is a dict with:
      - text: the combined text for that slide
      - metadata: slide_number, title, has_table, has_notes
    """
    chunks = []

    for slide in slides:
        parts = []

        # Title line
        if slide["title"]:
            parts.append(f"# {slide['title']}")

        # Body
        if slide["body_text"]:
            parts.append(slide["body_text"])

        # Tables rendered as text
        for tbl in slide["table_data"]:
            parts.append(table_to_text(tbl))

        # Speaker notes (optional — great for context)
        if include_notes and slide["notes"]:
            parts.append(f"[Notes] {slide['notes']}")

        text = "\n\n".join(parts).strip()
        if not text:
            continue

        chunks.append(
            {
                "text": text,
                "metadata": {
                    "slide_number": slide["slide_number"],
                    "title": slide["title"],
                    "has_table": len(slide["table_data"]) > 0,
                    "has_notes": bool(slide["notes"]),
                },
            }
        )

    return chunks


def build_slide_summaries(slides: list[dict]) -> list[str]:
    """
    Create short summaries for each slide, useful as context headers.

    Format: "Slide N — <Title>: <first 120 chars of body text>..."
    """
    summaries = []
    for slide in slides:
        body_preview = slide["body_text"][:120].replace("\n", " ")
        if len(slide["body_text"]) > 120:
            body_preview += "..."
        title = slide["title"] or "(untitled)"
        summaries.append(f"Slide {slide['slide_number']} — {title}: {body_preview}")
    return summaries


# ---------------------------------------------------------------------------
# Display helpers
# ---------------------------------------------------------------------------


def print_structured_slides(slides: list[dict]) -> None:
    """Pretty-print the structured slide data."""
    for slide in slides:
        print(f"\n{'=' * 60}")
        print(f"  SLIDE {slide['slide_number']}: {slide['title'] or '(no title)'}")
        print(f"{'=' * 60}")

        if slide["body_text"]:
            print("\n  Body text:")
            for line in slide["body_text"].split("\n"):
                print(f"    {line}")

        if slide["table_data"]:
            for t_idx, tbl in enumerate(slide["table_data"]):
                print(f"\n  Table {t_idx + 1} ({len(tbl)} rows x {len(tbl[0])} cols):")
                for row in tbl:
                    print(f"    {row}")

        if slide["notes"]:
            print("\n  Speaker notes:")
            for line in slide["notes"].split("\n"):
                print(f"    {line}")


def print_rag_chunks(chunks: list[dict]) -> None:
    """Print RAG chunks with metadata."""
    for i, chunk in enumerate(chunks, start=1):
        meta = chunk["metadata"]
        print(f"\n{'- ' * 30}")
        print(
            f"  Chunk {i}  |  Slide {meta['slide_number']}  |  "
            f"Title: {meta['title'] or 'N/A'}  |  "
            f"Table: {meta['has_table']}  |  Notes: {meta['has_notes']}"
        )
        print(f"{'- ' * 30}")
        # Truncate for display
        text = chunk["text"]
        if len(text) > 400:
            print(f"  {text[:400]}...")
        else:
            print(f"  {text}")


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    if not PPTX_PATH.exists():
        print(f"Sample file not found: {PPTX_PATH}")
        print("Run generate_samples.py first.")
        sys.exit(1)

    print("=" * 60)
    print("  Structured PPTX Extraction for RAG")
    print("=" * 60)
    print(f"  File: {PPTX_PATH.name}\n")

    # ── 1. Structured extraction ──────────────────────────────────────────
    slides = extract_structured_slides(PPTX_PATH)
    print("\n>>> STEP 1: Structured slide data")
    print_structured_slides(slides)

    # ── 2. Slide summaries ────────────────────────────────────────────────
    summaries = build_slide_summaries(slides)
    print(f"\n\n{'=' * 60}")
    print("  SLIDE SUMMARIES")
    print(f"{'=' * 60}")
    for s in summaries:
        print(f"  {s}")

    # ── 3. RAG-ready chunks (one per slide) ───────────────────────────────
    rag_chunks = slides_to_rag_chunks(slides, include_notes=True)
    print(f"\n\n{'=' * 60}")
    print("  RAG-READY CHUNKS (one per slide, notes included)")
    print(f"{'=' * 60}")
    print_rag_chunks(rag_chunks)

    # ── 4. Sentence-based chunking on full text ───────────────────────────
    full_text = "\n\n".join(chunk["text"] for chunk in rag_chunks)
    print(f"\n\n{'=' * 60}")
    print("  SENTENCE-BASED CHUNKING (merged text)")
    print(f"{'=' * 60}")
    sentence_chunks = chunk_by_sentences(full_text, sentences_per_chunk=5, overlap_sentences=1)
    preview_chunks(sentence_chunks, max_preview=5, max_chars=350)

    # ── 5. JSON output sample ─────────────────────────────────────────────
    print(f"\n\n{'=' * 60}")
    print("  JSON OUTPUT (first 2 slides)")
    print(f"{'=' * 60}")
    json_sample = json.dumps(rag_chunks[:2], indent=2, ensure_ascii=False)
    print(json_sample)
