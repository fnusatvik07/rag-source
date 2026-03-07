"""
Generate sample PPTX files for demonstrating RAG extraction techniques.

Creates two presentations:
  1. presentation.pptx  - 6-slide "Introduction to Machine Learning" deck
  2. data_presentation.pptx - 4-slide "Q4 Financial Review" deck

Run:
    uv run python unstructured_documents/03_pptx/sample_docs/generate_samples.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

SAMPLE_DIR = Path(__file__).resolve().parent  # .../sample_docs/


def _add_textbox(slide, left, top, width, height, text, font_size=14, bold=False):
    """Add a simple text box to a slide."""
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    return txbox


# ---------------------------------------------------------------------------
# Presentation 1 – Introduction to Machine Learning (6 slides)
# ---------------------------------------------------------------------------


def create_ml_presentation() -> Path:
    prs = Presentation()

    # ── Slide 1: Title slide ──────────────────────────────────────────────
    slide_layout = prs.slide_layouts[0]  # Title Slide layout
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Introduction to Machine Learning"
    slide.placeholders[1].text = "A practical overview of ML concepts, algorithms, and applications"

    # ── Slide 2: Bullet points – ML types ─────────────────────────────────
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Types of Machine Learning"
    body = slide.placeholders[1].text_frame
    body.text = "Machine learning is broadly categorized into three types:"
    bullets = [
        "Supervised Learning: learns from labeled data (classification, regression)",
        "Unsupervised Learning: finds hidden patterns in unlabeled data (clustering, dimensionality reduction)",
        "Reinforcement Learning: learns through trial-and-error interaction with an environment",
        "Semi-supervised and Self-supervised methods are gaining popularity as hybrid approaches",
    ]
    for bullet in bullets:
        p = body.add_paragraph()
        p.text = bullet
        p.level = 1

    # ── Slide 3: Table – Algorithm comparison ─────────────────────────────
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    _add_textbox(
        slide,
        Inches(0.5),
        Inches(0.3),
        Inches(9),
        Inches(0.6),
        "Comparison of ML Algorithms",
        font_size=24,
        bold=True,
    )

    rows, cols = 6, 4
    table_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.2), Inches(9), Inches(4))
    table = table_shape.table

    headers = ["Algorithm", "Type", "Use Case", "Complexity"]
    data = [
        ["Linear Regression", "Supervised", "Price prediction", "Low"],
        ["Random Forest", "Supervised", "Classification tasks", "Medium"],
        ["K-Means", "Unsupervised", "Customer segmentation", "Low"],
        ["Neural Network", "Supervised/Unsupervised", "Image recognition", "High"],
        ["Q-Learning", "Reinforcement", "Game playing", "High"],
    ]
    for j, h in enumerate(headers):
        table.cell(0, j).text = h
    for i, row in enumerate(data, start=1):
        for j, val in enumerate(row):
            table.cell(i, j).text = val

    # ── Slide 4: Content + speaker notes ──────────────────────────────────
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "The ML Pipeline"
    body = slide.placeholders[1].text_frame
    body.text = "Every ML project follows a common pipeline:"
    steps = [
        "Data Collection: gathering relevant datasets",
        "Data Preprocessing: cleaning, normalizing, feature engineering",
        "Model Selection: choosing the right algorithm for the task",
        "Training: fitting the model on training data",
        "Evaluation: measuring performance with metrics (accuracy, F1, RMSE)",
        "Deployment: serving the model in production",
    ]
    for step in steps:
        p = body.add_paragraph()
        p.text = step
        p.level = 1

    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = (
        "Speaker notes: Emphasize that data preprocessing typically consumes "
        "60-80% of project time. Model selection should be driven by the "
        "problem type (classification vs regression) and data characteristics "
        "(size, dimensionality, label availability). Mention AutoML tools "
        "like Google AutoML, H2O, and Auto-sklearn as alternatives for "
        "automated model selection."
    )

    # ── Slide 5: Grouped text boxes – key takeaways ───────────────────────
    slide_layout = prs.slide_layouts[5]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    _add_textbox(
        slide,
        Inches(0.5),
        Inches(0.3),
        Inches(9),
        Inches(0.6),
        "Key Takeaways",
        font_size=24,
        bold=True,
    )

    # Create a group shape containing three text boxes
    group = slide.shapes.add_group_shape()

    takeaways = [
        (
            "Data is King",
            "The quality and quantity of your data matters more than the algorithm you choose.",
        ),
        (
            "Start Simple",
            "Begin with simple models like linear regression before moving to complex architectures.",
        ),
        (
            "Iterate Fast",
            "Rapid experimentation and iteration lead to better results than perfecting a single approach.",
        ),
    ]
    box_width = Emu(Inches(2.8).emu)
    top = Emu(Inches(1.2).emu)

    for idx, (title, desc) in enumerate(takeaways):
        left = Emu(Inches(0.4 + idx * 3.1).emu)
        # Title box
        tb_title = group.shapes.add_textbox(left, top, box_width, Emu(Inches(0.6).emu))
        tb_title.text_frame.paragraphs[0].text = title
        tb_title.text_frame.paragraphs[0].font.size = Pt(18)
        tb_title.text_frame.paragraphs[0].font.bold = True
        tb_title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        # Description box
        tb_desc = group.shapes.add_textbox(left, Emu(top.emu + Inches(0.7).emu), box_width, Emu(Inches(2).emu))
        tb_desc.text_frame.word_wrap = True
        tb_desc.text_frame.paragraphs[0].text = desc
        tb_desc.text_frame.paragraphs[0].font.size = Pt(14)

    # ── Slide 6: Summary / conclusion ─────────────────────────────────────
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Summary and Next Steps"
    body = slide.placeholders[1].text_frame
    body.text = "In this presentation we covered:"
    conclusions = [
        "The three main types of machine learning",
        "Key algorithms and their trade-offs",
        "The end-to-end ML pipeline from data to deployment",
        "Practical tips: focus on data quality, start simple, iterate quickly",
    ]
    for c in conclusions:
        p = body.add_paragraph()
        p.text = c
        p.level = 1

    body.add_paragraph()  # blank line
    p = body.add_paragraph()
    p.text = "Next steps: explore hands-on tutorials, pick a dataset, and build your first model!"

    out_path = SAMPLE_DIR / "presentation.pptx"
    prs.save(str(out_path))
    print(f"Created: {out_path}")
    return out_path


# ---------------------------------------------------------------------------
# Presentation 2 – Q4 Financial Review (4 slides)
# ---------------------------------------------------------------------------


def create_data_presentation() -> Path:
    prs = Presentation()

    # ── Slide 1: Title ────────────────────────────────────────────────────
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Q4 Financial Review"
    slide.placeholders[1].text = "Fiscal Year 2024 — Confidential"

    # ── Slide 2: Revenue table ────────────────────────────────────────────
    slide_layout = prs.slide_layouts[5]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    _add_textbox(
        slide,
        Inches(0.5),
        Inches(0.3),
        Inches(9),
        Inches(0.6),
        "Quarterly Revenue Breakdown",
        font_size=24,
        bold=True,
    )

    rows, cols = 5, 4
    table_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.2), Inches(9), Inches(3.5))
    table = table_shape.table

    headers = ["Region", "Q1 ($M)", "Q2 ($M)", "Q3 ($M)"]
    data = [
        ["North America", "45.2", "48.7", "52.1"],
        ["Europe", "32.1", "34.5", "37.8"],
        ["Asia Pacific", "28.4", "31.2", "35.6"],
        ["Rest of World", "12.3", "13.1", "14.9"],
    ]
    for j, h in enumerate(headers):
        table.cell(0, j).text = h
    for i, row in enumerate(data, start=1):
        for j, val in enumerate(row):
            table.cell(i, j).text = val

    # ── Slide 3: Key metrics text boxes ───────────────────────────────────
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    _add_textbox(
        slide,
        Inches(0.5),
        Inches(0.3),
        Inches(9),
        Inches(0.6),
        "Key Financial Metrics",
        font_size=24,
        bold=True,
    )

    metrics = [
        ("Total Revenue", "$140.4M", "Up 12% YoY"),
        ("Gross Margin", "68.5%", "Improved from 65.2%"),
        ("Operating Income", "$32.1M", "22.9% operating margin"),
        ("Customer Count", "1,247", "Net new: 183 customers"),
    ]
    for idx, (label, value, note) in enumerate(metrics):
        row_top = Inches(1.3 + idx * 1.1)
        _add_textbox(
            slide,
            Inches(0.8),
            row_top,
            Inches(3),
            Inches(0.5),
            label,
            font_size=16,
            bold=True,
        )
        _add_textbox(
            slide,
            Inches(4.0),
            row_top,
            Inches(2),
            Inches(0.5),
            value,
            font_size=16,
            bold=False,
        )
        _add_textbox(
            slide,
            Inches(6.2),
            row_top,
            Inches(3),
            Inches(0.5),
            note,
            font_size=14,
            bold=False,
        )

    # ── Slide 4: Conclusion with bullets ──────────────────────────────────
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Outlook and Recommendations"
    body = slide.placeholders[1].text_frame
    body.text = "Based on Q4 performance:"
    points = [
        "Revenue growth is accelerating across all regions",
        "Asia Pacific shows the strongest growth trajectory at 25% YoY",
        "Gross margin improvements driven by operational efficiency",
        "Recommend increasing investment in APAC market expansion",
        "Target: $160M revenue for next fiscal year",
    ]
    for pt in points:
        p = body.add_paragraph()
        p.text = pt
        p.level = 1

    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = (
        "Discuss risk factors: currency fluctuations in APAC, "
        "potential regulatory changes in EU markets, and "
        "competitive pressure in North America from new entrants."
    )

    out_path = SAMPLE_DIR / "data_presentation.pptx"
    prs.save(str(out_path))
    print(f"Created: {out_path}")
    return out_path


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    print("Generating sample PPTX files ...\n")
    create_ml_presentation()
    create_data_presentation()
    print("\nDone! Sample files are ready in:", SAMPLE_DIR)
