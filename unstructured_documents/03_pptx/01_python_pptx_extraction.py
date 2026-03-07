"""
Basic PPTX text extraction using python-pptx.

Demonstrates how to:
  - Load a .pptx file and iterate over slides
  - Extract text from every shape type (text frames, tables, group shapes)
  - Retrieve speaker notes
  - Display results with slide numbers and shape types

Run:
    uv run python unstructured_documents/03_pptx/01_python_pptx_extraction.py
"""

import sys
from pathlib import Path

# --- shared chunking import ------------------------------------------------
sys.path.insert(0, str(Path(__file__).resolve().parents[2]))
from pptx import Presentation

from unstructured_documents.shared.chunking import (
    chunk_by_recursive_split,
    chunk_by_sentences,
    preview_chunks,
)

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
SAMPLE_DIR = Path(__file__).resolve().parent / "sample_docs"
PPTX_PATH = SAMPLE_DIR / "presentation.pptx"


# ---------------------------------------------------------------------------
# Extraction helpers
# ---------------------------------------------------------------------------


def extract_text_from_shape(shape) -> list[dict]:
    """
    Recursively extract text from a single shape.

    Returns a list of dicts: {"type": <shape_type>, "text": <extracted_text>}
    Handles text frames, tables, and group shapes (which may nest others).
    """
    results = []

    # --- Group shape: recurse into child shapes ---
    if shape.shape_type is not None and shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
        for child in shape.shapes:
            results.extend(extract_text_from_shape(child))
        return results

    # --- Table ---
    if shape.has_table:
        table = shape.table
        rows_text = []
        for row in table.rows:
            row_cells = [cell.text.strip() for cell in row.cells]
            rows_text.append(" | ".join(row_cells))
        table_text = "\n".join(rows_text)
        if table_text.strip():
            results.append({"type": "table", "text": table_text})
        return results

    # --- Text frame (covers placeholders, text boxes, etc.) ---
    if shape.has_text_frame:
        text = shape.text_frame.text.strip()
        if text:
            results.append({"type": "text_frame", "text": text})
        return results

    return results


def extract_notes(slide) -> str:
    """Return the speaker notes for a slide, or empty string."""
    if slide.has_notes_slide:
        notes_text = slide.notes_slide.notes_text_frame.text.strip()
        return notes_text
    return ""


def extract_table_data(slide) -> list[list[list[str]]]:
    """
    Return all tables on a slide as a list of 2-D lists.

    Each table is [[row0_cells], [row1_cells], ...].
    """
    tables = []
    for shape in slide.shapes:
        if shape.has_table:
            table_rows = []
            for row in shape.table.rows:
                table_rows.append([cell.text.strip() for cell in row.cells])
            tables.append(table_rows)
    return tables


# ---------------------------------------------------------------------------
# Full extraction
# ---------------------------------------------------------------------------


def extract_all_slides(pptx_path: Path) -> list[dict]:
    """
    Walk every slide and extract text, tables, and notes.

    Returns a list (one entry per slide) of:
      {
        "slide_number": int,
        "shapes": [{"type": str, "text": str}, ...],
        "notes": str,
      }
    """
    prs = Presentation(str(pptx_path))
    slides_data = []

    for idx, slide in enumerate(prs.slides, start=1):
        shape_extracts = []
        for shape in slide.shapes:
            shape_extracts.extend(extract_text_from_shape(shape))

        slides_data.append(
            {
                "slide_number": idx,
                "shapes": shape_extracts,
                "notes": extract_notes(slide),
            }
        )

    return slides_data


# ---------------------------------------------------------------------------
# Display
# ---------------------------------------------------------------------------


def print_extraction_results(slides_data: list[dict]) -> None:
    """Pretty-print the extraction results."""
    for slide in slides_data:
        print(f"\n{'=' * 60}")
        print(f"  SLIDE {slide['slide_number']}")
        print(f"{'=' * 60}")

        if not slide["shapes"]:
            print("  (no extractable text)")

        for item in slide["shapes"]:
            label = item["type"].upper()
            print(f"\n  [{label}]")
            for line in item["text"].split("\n"):
                print(f"    {line}")

        if slide["notes"]:
            print("\n  [SPEAKER NOTES]")
            for line in slide["notes"].split("\n"):
                print(f"    {line}")


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    if not PPTX_PATH.exists():
        print(f"Sample file not found: {PPTX_PATH}")
        print("Run generate_samples.py first.")
        sys.exit(1)

    print("=" * 60)
    print("  PPTX Text Extraction with python-pptx")
    print("=" * 60)
    print(f"  File: {PPTX_PATH.name}")

    # ── 1. Slide-by-slide extraction ──────────────────────────────────────
    slides_data = extract_all_slides(PPTX_PATH)
    print_extraction_results(slides_data)

    # ── 2. Table extraction detail ────────────────────────────────────────
    print(f"\n\n{'=' * 60}")
    print("  TABLE EXTRACTION DETAIL")
    print(f"{'=' * 60}")
    prs = Presentation(str(PPTX_PATH))
    for idx, slide in enumerate(prs.slides, start=1):
        tables = extract_table_data(slide)
        if tables:
            for t_idx, table in enumerate(tables):
                print(f"\n  Slide {idx} — Table {t_idx + 1} ({len(table)} rows x {len(table[0])} cols):")
                for row in table:
                    print(f"    {row}")

    # ── 3. Combine all text for chunking demo ─────────────────────────────
    all_text_parts = []
    for slide in slides_data:
        for item in slide["shapes"]:
            all_text_parts.append(item["text"])
        if slide["notes"]:
            all_text_parts.append(slide["notes"])

    full_text = "\n\n".join(all_text_parts)

    print(f"\n\n{'=' * 60}")
    print("  CHUNKING DEMO — Sentence-based")
    print(f"{'=' * 60}")
    sentence_chunks = chunk_by_sentences(full_text, sentences_per_chunk=4, overlap_sentences=1)
    preview_chunks(sentence_chunks, max_preview=4, max_chars=300)

    print(f"\n\n{'=' * 60}")
    print("  CHUNKING DEMO — Recursive split")
    print(f"{'=' * 60}")
    recursive_chunks = chunk_by_recursive_split(full_text, chunk_size=400)
    preview_chunks(recursive_chunks, max_preview=4, max_chars=300)
