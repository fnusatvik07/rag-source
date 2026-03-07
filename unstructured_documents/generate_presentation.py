"""
Generate a comprehensive PowerPoint presentation:
"Reading Various Data Sources for RAG"

Covers all 9 document types, internal format structures, extraction methods,
OCR deep dive, chunking strategies, and recommended approaches.

Run: uv run python unstructured_documents/generate_presentation.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Color scheme
# ---------------------------------------------------------------------------
DARK_BLUE = RGBColor(0x1F, 0x4E, 0x79)
MED_BLUE = RGBColor(0x2E, 0x75, 0xB6)
LIGHT_BLUE = RGBColor(0xD6, 0xE4, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
ACCENT_ORANGE = RGBColor(0xED, 0x7D, 0x31)
ACCENT_GREEN = RGBColor(0x70, 0xAD, 0x47)

OUTPUT_PATH = Path(__file__).parent / "RAG_Data_Sources_Presentation.pptx"


# ---------------------------------------------------------------------------
# Helper functions
# ---------------------------------------------------------------------------


def add_title_bar(slide, title_text: str):
    """Add a dark blue title bar at the top of a slide."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.15)


def add_section_divider(prs, section_title: str, subtitle: str = ""):
    """Add a full-slide section divider with dark background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = section_title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(20)
        p2.font.color.rgb = LIGHT_BLUE
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(20)

    return slide


def add_content_slide(prs, title: str, bullets: list[str], sub_bullets: dict[int, list[str]] | None = None):
    """Add a slide with title bar and bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_title_bar(slide, title)

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(12), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(17)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(8)
        p.space_after = Pt(4)
        p.level = 0

        # Add sub-bullets if provided
        if sub_bullets and i in sub_bullets:
            for sb in sub_bullets[i]:
                sp = tf.add_paragraph()
                sp.text = sb
                sp.font.size = Pt(15)
                sp.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
                sp.space_before = Pt(2)
                sp.level = 1

    return slide


def add_table_slide(prs, title: str, headers: list[str], rows: list[list[str]], intro: str = ""):
    """Add a slide with a formatted table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, title)

    top = Inches(1.5)
    if intro:
        txBox = slide.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(12), Inches(0.6))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = intro
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        top = Inches(2.0)

    n_rows = len(rows) + 1
    n_cols = len(headers)
    col_width = min(Inches(12) // n_cols, Inches(3))
    table_width = col_width * n_cols

    table_shape = slide.shapes.add_table(n_rows, n_cols, Inches(0.6), top, table_width, Inches(0.4) * n_rows)
    table = table_shape.table

    # Header row
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(13)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
            paragraph.alignment = PP_ALIGN.CENTER

    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if i % 2 == 0 else LIGHT_GRAY
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.font.color.rgb = DARK_GRAY

    return slide


def add_two_column_slide(
    prs,
    title: str,
    left_title: str,
    left_bullets: list[str],
    right_title: str,
    right_bullets: list[str],
):
    """Add a slide with two columns."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, title)

    for col_x, col_title, col_bullets in [
        (Inches(0.5), left_title, left_bullets),
        (Inches(6.8), right_title, right_bullets),
    ]:
        # Column header
        hdr = slide.shapes.add_textbox(col_x, Inches(1.4), Inches(5.8), Inches(0.5))
        hp = hdr.text_frame.paragraphs[0]
        hp.text = col_title
        hp.font.size = Pt(20)
        hp.font.bold = True
        hp.font.color.rgb = MED_BLUE

        # Column content
        body = slide.shapes.add_textbox(col_x, Inches(2.0), Inches(5.8), Inches(5))
        tf = body.text_frame
        tf.word_wrap = True
        for i, b in enumerate(col_bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = b
            p.font.size = Pt(15)
            p.font.color.rgb = DARK_GRAY
            p.space_before = Pt(6)

    return slide


# =========================================================================
# BUILD THE PRESENTATION
# =========================================================================


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # =================================================================
    # SLIDE 1: Title
    # =================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11), Inches(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Reading Various Data Sources for RAG"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "A Comprehensive Guide to Parsing Unstructured Documents"
    p2.font.size = Pt(24)
    p2.font.color.rgb = LIGHT_BLUE
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(20)

    p3 = tf.add_paragraph()
    p3.text = "PDF | DOCX | PPTX | HTML | Spreadsheets | Images/OCR | Email | Markdown | EPUB"
    p3.font.size = Pt(18)
    p3.font.color.rgb = RGBColor(0x90, 0xB0, 0xD0)
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(40)

    # =================================================================
    # SLIDE 2: What is RAG?
    # =================================================================
    add_content_slide(
        prs,
        "What is RAG? (Retrieval-Augmented Generation)",
        [
            "RAG combines a Large Language Model (LLM) with an external knowledge base",
            "Instead of relying solely on training data, the LLM retrieves relevant documents at query time",
            "The quality of RAG depends entirely on the quality of your document parsing",
            "Garbage in = Garbage out: poorly parsed documents lead to irrelevant retrieval",
            "This presentation covers HOW to extract text from every major document type",
        ],
        {
            0: ["User asks a question → System retrieves relevant chunks → LLM generates answer using those chunks"],
            3: ["If your PDF text is garbled, table data is mangled, or images aren't OCR'd — the LLM can't help"],
        },
    )

    # =================================================================
    # SLIDE 3: The Document Parsing Challenge
    # =================================================================
    add_content_slide(
        prs,
        "The Document Parsing Challenge",
        [
            'Why can\'t we just "read" a file?',
            "Each file format stores data in fundamentally different ways internally",
            "A PDF doesn't contain 'paragraphs' — it contains coordinates and character codes",
            "A DOCX isn't a text file — it's a ZIP archive containing XML",
            "An image has NO text at all — just millions of colored pixels",
            "Tables, headers, footers, and multi-column layouts add complexity",
            "The same logical content looks completely different across formats",
            "Different methods give different results — choosing the right one matters",
        ],
    )

    # =================================================================
    # SLIDE 4: Overview of 9 Data Sources
    # =================================================================
    add_table_slide(
        prs,
        "9 Unstructured Data Sources for RAG",
        ["Type", "Format", "Key Challenge", "# Methods"],
        [
            [
                "PDF",
                "Binary / coordinate-based",
                "Text encoding, tables, scanned docs",
                "6",
            ],
            [
                "Word (DOCX)",
                "ZIP of XML files",
                "Styles, heading hierarchy, tables",
                "3",
            ],
            [
                "PowerPoint (PPTX)",
                "ZIP of XML slides",
                "Sparse text, visual content, notes",
                "2",
            ],
            [
                "HTML / Web",
                "DOM tree markup",
                "Boilerplate removal (nav, ads, footer)",
                "3",
            ],
            [
                "Spreadsheets",
                "Cell-based (XLSX/CSV)",
                "Meaning from position, not text",
                "3",
            ],
            ["Images", "Pixel grid (PNG/JPG)", "No inherent text — needs OCR", "2"],
            [
                "Email (EML)",
                "MIME multipart",
                "Encoded parts, HTML bodies, attachments",
                "2",
            ],
            [
                "Markdown / Text",
                "Plain text",
                "No structure — chunking is the challenge",
                "3",
            ],
            [
                "EPUB (Ebooks)",
                "ZIP of XHTML chapters",
                "Chapter extraction, HTML inside",
                "2",
            ],
        ],
        intro="Each document type stores data differently and requires specialized extraction approaches.",
    )

    # =================================================================
    # SLIDE 5: Universal RAG Pipeline
    # =================================================================
    add_content_slide(
        prs,
        "The Universal RAG Document Pipeline",
        [
            "Every document type follows the same high-level pipeline:",
            "",
            "1. SOURCE DOCUMENT  →  Your PDF, DOCX, HTML, image, etc.",
            "2. PARSE / EXTRACT  →  Pull out raw text, tables, metadata (THIS presentation focuses here)",
            "3. CLEAN            →  Remove noise, normalize formatting, strip boilerplate",
            "4. CHUNK            →  Split into retrieval-sized pieces (200-500 tokens typical)",
            "5. EMBED            →  Convert each chunk to a vector using an embedding model",
            "6. STORE            →  Save vectors + metadata in a vector database (Pinecone, Chroma, etc.)",
            "",
            "Steps 2-4 (Parse, Clean, Chunk) are where most RAG quality is won or lost",
            "This presentation covers steps 2-4 in depth for every document type",
        ],
    )

    # =================================================================
    # SLIDE 6: Why Different Approaches
    # =================================================================
    add_two_column_slide(
        prs,
        "Why Different Formats Need Different Approaches",
        "What You See (Visual)",
        [
            "A heading in bold, large font",
            "A paragraph of body text",
            "A table with rows and columns",
            "A bulleted list of items",
            "An image with caption",
            "",
            "All of these LOOK similar across formats...",
        ],
        "What The Computer Sees (Internal)",
        [
            "PDF: BT /F1 24 Tf 72 720 Td (Heading) Tj ET",
            'DOCX: <w:pStyle w:val="Heading1"/>',
            "PPTX: <a:t>Heading</a:t> inside <p:sp>",
            "HTML: <h1>Heading</h1>",
            "Image: [0,0,0,255,255,255,0,0,0...] pixels",
            "",
            "...but are stored completely differently!",
        ],
    )

    # =================================================================
    # SLIDE 7: How Data Lives Inside Files
    # =================================================================
    add_content_slide(
        prs,
        "How Data Actually Lives Inside Files",
        [
            (
                "PDF: Binary content streams with character codes + (x,y) coordinates."
                " The text 'Hello' might be stored as individual characters placed at"
                " specific pixel positions. No concept of 'paragraphs'."
            ),
            (
                "DOCX: A ZIP file containing XML. Unzip a .docx and you'll find"
                " word/document.xml with <w:p> (paragraph) and <w:r> (run) tags."
            ),
            (
                "PPTX: A ZIP file of XML slides. Each slide has shapes (<p:sp>)"
                " containing text frames. Text isn't linear — it's in positioned boxes."
            ),
            "HTML: A tree of nested tags. Actual content is mixed with navigation, ads, footers, and JavaScript.",
            (
                "Spreadsheets: Cells addressed by (row, col). 'Salary: $125,000' only"
                " makes sense because the header is in row 1."
            ),
            (
                "Images: Just a grid of pixel colors (R, G, B values). The letter 'A'"
                " is a pattern of dark pixels on a light background — the computer has"
                " zero knowledge it's text."
            ),
            "Email: MIME-encoded parts, potentially Base64, with headers and multipart boundaries.",
        ],
    )

    # =================================================================
    # SLIDE 8: Key Concepts
    # =================================================================
    add_content_slide(
        prs,
        "Key Concepts for This Presentation",
        [
            "Parsing: Reading the raw file format and extracting its contents",
            "Extraction: Pulling meaningful text, tables, and metadata from parsed content",
            "Chunking: Splitting extracted text into smaller pieces sized for embedding models",
            "Tokens: Sub-word units that LLMs process (~1 token ≈ 4 characters in English)",
            "Embeddings: Dense vector representations of text chunks (e.g., 1536 dimensions)",
            (
                "Overlap: Including some text from the previous chunk in the next one"
                " to avoid losing context at boundaries"
            ),
            "Metadata: Information ABOUT the chunk (source file, page number, heading, date) used for filtering",
            "Boilerplate: Non-content elements (navigation, footers, ads) that add noise to embeddings",
        ],
    )

    # =================================================================
    # SECTION 2: File Format Internals
    # =================================================================
    add_section_divider(prs, "Section 2", "How Data Lives Inside Each Format")

    # SLIDE 9: PDF Internals
    add_content_slide(
        prs,
        "PDF Internals: Not What You Think",
        [
            "PDF = Portable Document Format. Designed for DISPLAY, not data extraction.",
            "Text is NOT stored as paragraphs — it's individual characters at (x, y) coordinates",
            "Example: The word 'Hello' might be stored as:",
            "    BT /F1 12 Tf 72 720 Td (H) Tj 7.2 0 Td (e) Tj 5.4 0 Td (l) Tj ...",
            "Fonts can use custom encoding — character code 65 doesn't always mean 'A'",
            "Images in PDFs are embedded binary streams — text inside images is invisible to parsers",
            "Tables are NOT tables — they're just text characters aligned in a grid visually",
            "Multi-column layouts: text flows down column 1, then column 2 — parsers may interleave them",
            "This is why PDF parsing has 6 different methods — each handles these challenges differently",
        ],
    )

    # SLIDE 10: DOCX Internals
    add_content_slide(
        prs,
        "DOCX Internals: ZIP of XML Files",
        [
            "A .docx file is literally a ZIP archive. Rename it to .zip and extract it!",
            "Inside you'll find: word/document.xml, word/styles.xml, word/media/ (images), etc.",
            "document.xml structure: <w:body> → <w:p> (paragraph) → <w:r> (run) → <w:t> (text)",
            "Styles define Heading 1, Heading 2, Normal, etc. — essential for structured extraction",
            "Each 'run' can have different formatting (bold, italic, font) within the same paragraph",
            "Tables are <w:tbl> → <w:tr> (row) → <w:tc> (cell) — well-structured for extraction",
            "Advantage over PDF: paragraphs and headings are explicitly marked in the XML",
            "This is why DOCX parsing is generally easier and more reliable than PDF parsing",
        ],
    )

    # SLIDE 11: PPTX Internals
    add_content_slide(
        prs,
        "PPTX Internals: Slides, Shapes, and Text Frames",
        [
            "A .pptx file is also a ZIP archive, similar to DOCX",
            "Inside: ppt/slides/slide1.xml, slide2.xml, etc. + notesSlides/ for speaker notes",
            "Each slide contains shapes (<p:sp>) — text boxes, tables, images, group shapes",
            "Text is inside text frames: <p:txBody> → <a:p> (paragraph) → <a:r> (run) → <a:t> (text)",
            "Title and content are in 'placeholders' — but not all shapes are placeholders",
            "Group shapes nest other shapes inside them — extractors must recurse",
            "Speaker notes live in separate notesSlide XML files — often contain the BEST text for RAG",
            "Challenge: Slides are visual — text is sparse and context depends on visual layout",
        ],
    )

    # SLIDE 12: HTML Internals
    add_content_slide(
        prs,
        "HTML Internals: DOM Tree with Boilerplate",
        [
            "HTML is a tree of nested tags: <html> → <body> → <div> → <p> → text",
            "Semantic tags: <h1>, <article>, <nav>, <footer>, <table> — meaningful for extraction",
            "The main content might be 30% of the HTML — the rest is navigation, ads, scripts, CSS",
            "Tables in HTML are well-structured: <table> → <tr> → <td> — easier to parse than PDF tables",
            "Links (<a href>), images (<img>), emphasis (<strong>, <em>) can be preserved or stripped",
            "JavaScript-rendered content (SPAs) won't be visible to HTML parsers — need Selenium/Playwright",
            "Encoding issues: pages may use UTF-8, Latin-1, or other character encodings",
            "The key challenge is separating the article/main content from everything else",
        ],
    )

    # SLIDE 13: Spreadsheet Internals
    add_content_slide(
        prs,
        "Spreadsheet Internals: Cells, Types, and Relationships",
        [
            "XLSX: ZIP of XML files. Each sheet is a separate XML file with rows and cells",
            "Each cell has: a position (A1, B3), a type (string, number, date, formula), and a value",
            "Meaning comes from POSITION — cell B5 only makes sense in context of its header in B1",
            "Merged cells: A header spanning columns A-D creates NULL values in B, C, D",
            "Formulas: Cell D5 might contain '=SUM(A5:C5)' — you want the calculated value, not the formula",
            "Multiple sheets: One workbook can have many sheets, each a separate 'document'",
            "CSV is simpler: plain text with delimiter-separated values, but loses types and formatting",
            "Challenge for RAG: 'Salary: $125,000' is meaningful, but '$125,000' alone is not",
        ],
    )

    # SLIDE 14: Image Internals
    add_content_slide(
        prs,
        "Image Internals: Pixels — No Text At All",
        [
            "An image is a 2D grid of pixels. Each pixel has color values (Red, Green, Blue: 0-255)",
            "Example: A 800x600 image = 480,000 pixels = 1,440,000 color values",
            "The letter 'A' in an image is just a pattern: dark pixels on a lighter background",
            "The computer has ZERO knowledge that these pixels form the letter 'A'",
            "Image formats: PNG (lossless), JPEG (lossy compression), TIFF (high quality scans)",
            "Resolution matters: 72 DPI (screen) vs 300 DPI (print quality) — higher DPI = better OCR",
            "Grayscale images: Each pixel is just one value (0=black, 255=white) instead of RGB",
            "To extract text from images, we need OCR — Optical Character Recognition",
            "OCR converts pixel patterns back into characters — a complex, error-prone process",
        ],
    )

    # SLIDE 15: Email Internals
    add_content_slide(
        prs,
        "Email Internals: MIME Multipart Structure",
        [
            "Email format (EML) uses MIME (Multipurpose Internet Mail Extensions)",
            "Headers: From, To, CC, Subject, Date, Message-ID — all structured key-value pairs",
            "Body can be: plain text, HTML, or BOTH (multipart/alternative)",
            "Attachments are Base64-encoded binary data inside multipart/mixed boundaries",
            "Example MIME structure:",
            "    Content-Type: multipart/mixed; boundary='----=_Part_1234'",
            "    ------=_Part_1234",
            "    Content-Type: text/plain → the body text",
            "    ------=_Part_1234",
            "    Content-Type: application/pdf; name='report.pdf' → attachment",
            "Headers can use RFC 2047 encoding for non-ASCII characters",
        ],
    )

    # SLIDE 16: Markdown/Text Internals
    add_content_slide(
        prs,
        "Markdown & Plain Text: The Simplest Format",
        [
            "Plain text: Just characters. No formatting metadata at all.",
            "Markdown adds lightweight conventions: # Heading, **bold**, - bullet, [link](url)",
            "These are NOT parsed by the file format — they're just text patterns we interpret",
            "No binary data, no XML, no encoding tricks — just UTF-8 text",
            "Heading structure: # = H1, ## = H2, ### = H3 — perfect natural chunk boundaries",
            "Code blocks: Fenced with ``` — should be extracted separately for technical RAG",
            "Since there's no 'parsing' challenge, the main challenge is CHUNKING strategy",
            "Plain text without headings: You must rely on paragraph breaks, sentences, or fixed sizes",
        ],
    )

    # SLIDE 17: EPUB Internals
    add_content_slide(
        prs,
        "EPUB Internals: A Book in a ZIP",
        [
            "EPUB = Electronic Publication. It's a ZIP file containing XHTML chapters",
            "Structure: META-INF/container.xml → content.opf (spine/manifest) → chapter files",
            "The 'spine' defines reading order: chapter1.xhtml → chapter2.xhtml → ...",
            "Each chapter is a complete XHTML document with HTML tags inside",
            "Metadata is rich: title, author, publisher, language, ISBN, description",
            "Table of Contents (TOC) is in toc.ncx or nav.xhtml — maps to chapter files",
            "Great for RAG: natural chapter boundaries provide excellent chunk boundaries",
            "Parsing: Read the EPUB with ebooklib, then parse each chapter's HTML with BeautifulSoup",
        ],
    )

    # =================================================================
    # SECTION 3: PDF Deep Dive
    # =================================================================
    add_section_divider(prs, "Section 3", "PDF Deep Dive — The Most Complex Format")

    # SLIDE 18: Why PDF is Hard
    add_table_slide(
        prs,
        "Why PDF Parsing is Hard — 7 Key Challenges",
        ["Challenge", "Description", "Impact on RAG"],
        [
            [
                "Text Encoding",
                "Characters mapped through font-specific encoding tables",
                "Garbled or missing text",
            ],
            [
                "Layout vs Content",
                "Text placed by (x,y) coordinates, not logical order",
                "Jumbled reading order",
            ],
            [
                "Multi-Column",
                "Text flows down column 1, then column 2",
                "Columns get interleaved",
            ],
            [
                "Tables",
                "No table structure — just aligned characters",
                "Tables become garbled text",
            ],
            [
                "Scanned PDFs",
                "Pages are images, not text at all",
                "Zero text extracted without OCR",
            ],
            [
                "Headers/Footers",
                "Repeated on every page with no marker",
                "Noise in every chunk",
            ],
            [
                "Embedded Fonts",
                "Custom fonts with non-standard character maps",
                "Special characters lost",
            ],
        ],
    )

    # SLIDE 19: pypdf
    add_two_column_slide(
        prs,
        "PDF Method 1: pypdf — Basic Text Extraction",
        "How It Works",
        [
            "Reads PDF structure and extracts text streams",
            "Decodes character codes using font maps",
            "Concatenates text in the order found in the PDF",
            "Also extracts metadata (title, author, dates)",
            "Pure Python — no system dependencies needed",
        ],
        "Pros & Cons",
        [
            "PROS:",
            "  + Fast (6ms per file in benchmarks)",
            "  + Lightweight, pure Python",
            "  + Good metadata extraction",
            "  + Active development, well maintained",
            "",
            "CONS:",
            "  - No table detection",
            "  - No layout awareness",
            "  - Can jumble multi-column text",
            "  - No OCR support",
        ],
    )

    # SLIDE 20: pdfplumber
    add_two_column_slide(
        prs,
        "PDF Method 2: pdfplumber — Tables & Layout",
        "How It Works",
        [
            "Builds on pdfminer for text extraction",
            "Analyzes character positions to detect tables",
            "Offers 'layout' mode preserving spatial positioning",
            "Provides word-level and character-level data",
            "Returns tables as structured Python lists",
        ],
        "Pros & Cons",
        [
            "PROS:",
            "  + Best table detection and extraction",
            "  + Layout-aware mode for multi-column",
            "  + Character-level detail (fonts, sizes)",
            "  + Word bounding boxes available",
            "",
            "CONS:",
            "  - Slower than pypdf (~95ms per file)",
            "  - No OCR support",
            "  - More memory usage",
            "  - Complex API for simple use cases",
        ],
    )

    # SLIDE 21: PyMuPDF
    add_two_column_slide(
        prs,
        "PDF Method 3: PyMuPDF (fitz) — Speed & Power",
        "How It Works",
        [
            "C-based engine (MuPDF) with Python bindings",
            "Extracts text as 'blocks' with bounding boxes",
            "Can detect text properties (font, size, bold)",
            "Sort mode reorders blocks by reading position",
            "Structured dict extraction for font analysis",
        ],
        "Pros & Cons",
        [
            "PROS:",
            "  + Fastest library (5.5ms per file)",
            "  + Block-level position data",
            "  + Font size/name detection → heading detection",
            "  + Good for batch processing at scale",
            "",
            "CONS:",
            "  - C dependency (binary install)",
            "  - Table extraction is basic",
            "  - Sorted mode can be slower (~25ms)",
            "  - License considerations (AGPL)",
        ],
    )

    # SLIDE 22: PDF Table Extraction
    add_content_slide(
        prs,
        "PDF Table Extraction — Converting Tables for RAG",
        [
            "Tables in PDFs are NOT real tables — they're characters aligned visually",
            "pdfplumber detects tables by analyzing character positions and line segments",
            "",
            "Key strategy: Convert extracted tables into natural language descriptions:",
            "  BEFORE (raw table): ['Product', 'Price'] → ['Widget', '$29.99']",
            "  AFTER (RAG-ready):  'Product: Widget; Price: $29.99; Category: Electronics'",
            "",
            "Why? Because 'Widget $29.99' without context means nothing to an embedding model",
            "But 'Product Widget costs $29.99' is a complete, searchable statement",
            "",
            "Alternative formats: Markdown tables, CSV strings, JSON — all work but natural language is best",
            "Always include column headers with each row — they provide the semantic context",
        ],
    )

    # SLIDE 23: OCR for Scanned PDFs
    add_content_slide(
        prs,
        "OCR for Scanned PDFs — When Text Extraction Returns Empty",
        [
            "Scanned PDFs are images wrapped in PDF format — no extractable text layer",
            "How to detect: If pypdf/pdfplumber return empty or near-empty text → likely scanned",
            "",
            "The OCR Pipeline for PDFs:",
            "  1. pdf2image: Convert each PDF page to a high-resolution image (300 DPI recommended)",
            "  2. Preprocessing: Convert to grayscale, apply threshold, sharpen, denoise",
            "  3. pytesseract: Run Tesseract OCR on each image → get text",
            "  4. Post-process: Clean up OCR artifacts, fix common errors",
            "",
            "Performance: OCR is 10-100x slower than direct text extraction",
            "Accuracy: ~95% for clean scans, drops significantly for low-quality or handwritten docs",
            "Rule: Always try direct extraction first. Only use OCR when it returns empty/garbage",
        ],
    )

    # SLIDE 24: PDF Method Comparison
    add_table_slide(
        prs,
        "PDF Method Comparison & Decision Guide",
        ["Method", "Speed", "Tables", "Layout", "OCR", "Best For"],
        [
            [
                "pypdf",
                "6ms (fastest)",
                "No",
                "No",
                "No",
                "Simple text, metadata, bulk processing",
            ],
            [
                "pdfplumber",
                "95ms",
                "Excellent",
                "Yes (layout mode)",
                "No",
                "Tables, forms, structured data",
            ],
            [
                "PyMuPDF",
                "5.5ms (fastest)",
                "Basic",
                "Block positions",
                "No",
                "Speed-critical, font analysis",
            ],
            [
                "pytesseract",
                "2-10 seconds",
                "No",
                "hOCR available",
                "Yes",
                "Scanned/image PDFs only",
            ],
        ],
        intro="Decision: Try pypdf/PyMuPDF first → Use pdfplumber for tables → Use OCR only for scanned PDFs",
    )

    # =================================================================
    # SECTION 4: DOCX Deep Dive
    # =================================================================
    add_section_divider(prs, "Section 4", "DOCX Deep Dive — ZIP of XML")

    # SLIDE 25: DOCX Structure
    add_content_slide(
        prs,
        "DOCX Structure: Paragraphs, Runs, and Styles",
        [
            "Every DOCX is a hierarchy: Document → Paragraphs → Runs → Text",
            "",
            "Paragraph (<w:p>): A block of text with one style (Heading 1, Normal, List Bullet, etc.)",
            "Run (<w:r>): A span within a paragraph with consistent formatting (bold, italic, font)",
            "Text (<w:t>): The actual characters within a run",
            "",
            'Example: "This is bold and normal text" is stored as:',
            "  Paragraph [style=Normal]",
            '    Run [bold=True]: "This is bold"',
            '    Run [bold=False]: " and normal text"',
            "",
            "Style names are KEY for RAG: They tell you what's a heading vs body vs list item",
            "This heading hierarchy enables heading-aware chunking — the best strategy for DOCX",
        ],
    )

    # SLIDE 26: python-docx
    add_content_slide(
        prs,
        "python-docx: Full Structural Extraction",
        [
            "python-docx gives you direct access to the DOCX XML structure via Python objects",
            "",
            "Key capabilities:",
            "  doc.paragraphs → iterate all paragraphs with their style names",
            "  paragraph.style.name → 'Heading 1', 'Normal', 'List Bullet', etc.",
            "  paragraph.runs → access individual formatted text spans",
            "  doc.tables → extract all tables with row/cell access",
            "",
            "Extraction strategy: Build a heading hierarchy",
            "  For each paragraph: if style is 'Heading N' → start new section",
            "  Group all body paragraphs under their parent heading",
            "  Each section becomes one RAG chunk with the heading as metadata",
            "",
            "Best for: When you need the document's structure preserved for heading-aware chunking",
        ],
    )

    # SLIDE 27: mammoth & docx2txt
    add_two_column_slide(
        prs,
        "mammoth vs docx2txt — Two Simpler Approaches",
        "mammoth (Convert to Markdown)",
        [
            "Converts DOCX to clean markdown or HTML",
            "Maps styles to semantic tags automatically:",
            "  Heading 1 → # Heading 1",
            "  Bold → **bold**",
            "  List Bullet → - item",
            "",
            "Best for: Quick markdown output that's",
            "ready for heading-aware chunking",
            "",
            "Output is immediately RAG-friendly!",
        ],
        "docx2txt (Plain Text Dump)",
        [
            "One-line extraction: docx2txt.process(path)",
            "Returns plain text only — no styles preserved",
            "",
            "What's LOST:",
            "  - Heading levels (all become plain text)",
            "  - Bold/italic markers",
            "  - Bullet point indicators",
            "  - Hyperlink URLs",
            "",
            "Best for: When you just need raw text",
            "and will use character/sentence chunking",
        ],
    )

    # SLIDE 28: DOCX Comparison
    add_table_slide(
        prs,
        "DOCX Method Comparison & Recommended Approach",
        ["Method", "Preserves Styles", "Tables", "Markdown Output", "Best For"],
        [
            [
                "python-docx",
                "Full (name, bold, etc.)",
                "Yes (structured)",
                "Manual conversion",
                "Full structural extraction",
            ],
            [
                "mammoth",
                "Semantic mapping",
                "Basic",
                "Native markdown",
                "Quick markdown, heading chunking",
            ],
            [
                "docx2txt",
                "None (plain text)",
                "Cell text only",
                "No",
                "Simple text dump, minimal deps",
            ],
        ],
        intro="Recommendation: python-docx for full control → mammoth for quick markdown → docx2txt for simplicity",
    )

    # =================================================================
    # SECTION 5: PPTX Deep Dive
    # =================================================================
    add_section_divider(prs, "Section 5", "PPTX Deep Dive — Sparse Text in Visual Containers")

    # SLIDE 29: Why PPTX is Tricky
    add_content_slide(
        prs,
        "Why PPTX is Tricky for RAG",
        [
            "Presentations are inherently VISUAL — designed for human viewing, not text extraction",
            "",
            "7 challenges unique to PPTX:",
            "  1. Sparse text: Each slide might have only 5-20 words of actual content",
            "  2. Visual context: A chart or diagram conveys meaning that text extraction misses",
            "  3. Fragmented text: Content is scattered across multiple shapes per slide",
            "  4. Group shapes: Nested shapes require recursive extraction",
            "  5. Speaker notes: The BEST text is often in notes, not on the slide",
            "  6. Tables in slides: Different from DOCX tables, stored as shape sub-elements",
            "  7. Slide ordering: Content flow depends on visual position, not XML order",
            "",
            "Key insight: Speaker notes often contain the full explanation for each slide",
        ],
    )

    # SLIDE 30: PPTX Extraction
    add_two_column_slide(
        prs,
        "PPTX Extraction: Basic vs Structured Approach",
        "Basic Extraction",
        [
            "Walk all slides → all shapes → extract text",
            "Handles: text frames, tables, group shapes",
            "Returns flat text per slide",
            "",
            "Process:",
            "  for slide in presentation.slides:",
            "    for shape in slide.shapes:",
            "      if shape.has_text_frame:",
            "        text = shape.text_frame.text",
            "      if shape.has_table:",
            "        # extract table cells",
        ],
        "Structured Extraction (Recommended)",
        [
            "Parse each slide into a structured dict:",
            "  {slide_number, title, body_text,",
            "   table_data, notes}",
            "",
            "Key improvements:",
            "  - Separate title from body text",
            "  - Extract speaker notes",
            "  - Convert tables to 'Header: Value' format",
            "  - Build slide summaries as context",
            "",
            "Creates richer, more meaningful chunks",
            "with metadata for better retrieval",
        ],
    )

    # SLIDE 31: PPTX Strategy
    add_content_slide(
        prs,
        "PPTX Strategy: Slide-Per-Chunk with Metadata",
        [
            "Recommended approach: One chunk per slide with rich metadata",
            "",
            "Each chunk should contain:",
            "  - Slide title as a heading",
            "  - All body text from the slide",
            "  - Table data converted to natural language",
            "  - Speaker notes (these are gold for RAG!)",
            "",
            "Metadata for each chunk:",
            "  - slide_number, title, has_table, has_notes",
            "",
            "Why per-slide chunking works:",
            "  - Slides are natural topic boundaries (the presenter chose them)",
            "  - Each slide covers one concept or point",
            "  - Speaker notes provide the context that sparse slide text lacks",
            "  - Metadata enables filtering ('find slides with tables')",
        ],
    )

    # =================================================================
    # SECTION 6: HTML Deep Dive
    # =================================================================
    add_section_divider(prs, "Section 6", "HTML Deep Dive — Separating Content from Noise")

    # SLIDE 32: Boilerplate Problem
    add_content_slide(
        prs,
        "The HTML Boilerplate Problem",
        [
            "A typical web page is 70% boilerplate, 30% actual content",
            "",
            "What's boilerplate?",
            "  - Navigation menus (<nav>): Home | About | Contact | Blog",
            "  - Sidebars (<aside>): Related articles, newsletter signup",
            "  - Footers (<footer>): Copyright, privacy policy, terms",
            "  - Advertisements: Injected scripts and iframes",
            "  - Cookie banners, popup overlays",
            "  - JavaScript, CSS, metadata tags",
            "",
            "If you embed this boilerplate: Every chunk will contain 'Home | About | Contact'",
            "This dilutes the semantic meaning and hurts retrieval quality",
            "",
            "Solution: Remove boilerplate BEFORE chunking. Use trafilatura (automatic) or BeautifulSoup (manual)",
        ],
    )

    # SLIDE 33: BeautifulSoup
    add_content_slide(
        prs,
        "BeautifulSoup: Fine-Grained HTML Control",
        [
            "BeautifulSoup parses HTML into a navigable tree of Python objects",
            "",
            "Key operations for RAG extraction:",
            "  soup.find('article') → find the main content container",
            "  soup.find_all(['nav', 'footer']).decompose() → remove boilerplate",
            "  soup.find_all('h2') → extract all headings",
            "  soup.find_all('table') → extract all tables separately",
            "  soup.get_text(separator='\\n') → get clean text from any element",
            "",
            "Best practice: Target the <article> or <main> tag, then extract from there",
            "This skips navigation, footer, and sidebar automatically",
            "",
            "Best for: Custom scrapers for known site structures, API documentation, product pages",
        ],
    )

    # SLIDE 34: html2text & Trafilatura
    add_two_column_slide(
        prs,
        "html2text vs Trafilatura",
        "html2text (HTML → Markdown)",
        [
            "Converts ALL HTML to markdown text",
            "No boilerplate removal (gets everything)",
            "Great markdown output with headings (#)",
            "",
            "RAG-optimized settings:",
            "  converter.ignore_links = True",
            "  converter.ignore_images = True",
            "  converter.body_width = 0",
            "",
            "Best for: Well-structured HTML where",
            "you want heading-aware chunking",
        ],
        "Trafilatura (Smart Extraction)",
        [
            "Designed specifically for web articles",
            "AUTOMATICALLY removes boilerplate",
            "Extracts just the main content",
            "Also extracts metadata (author, date)",
            "",
            "Supports output formats:",
            "  - Plain text (default)",
            "  - JSON (with metadata)",
            "  - XML (structural markup)",
            "",
            "Best for: Web scraping, diverse sites,",
            "articles/blogs with unknown structure",
        ],
    )

    # SLIDE 35: HTML Comparison
    add_table_slide(
        prs,
        "HTML Method Comparison & Decision Guide",
        ["Method", "Boilerplate Removal", "Tables", "Markdown Output", "Best For"],
        [
            [
                "BeautifulSoup",
                "Manual (you control)",
                "Yes (manual)",
                "No",
                "Custom scrapers, known structures",
            ],
            [
                "html2text",
                "No (converts all)",
                "Auto",
                "Yes (native)",
                "Quick markdown, heading chunking",
            ],
            ["Trafilatura", "Automatic", "Yes", "No", "Web articles, unknown sites"],
        ],
        intro="Decision: Trafilatura for web scraping → html2text for markdown → BeautifulSoup for custom control",
    )

    # =================================================================
    # SECTION 7: Spreadsheets
    # =================================================================
    add_section_divider(prs, "Section 7", "Spreadsheets Deep Dive — When Meaning Comes from Position")

    # SLIDE 36: Why Tables are Hard
    add_content_slide(
        prs,
        "Why Tabular Data is Hard for RAG",
        [
            "In text documents, meaning is in the words: 'The revenue was $2.4M'",
            "In spreadsheets, meaning comes from POSITION: Cell B5 = $2,400,000",
            "",
            "Without the header row, '$2,400,000' is meaningless",
            "Without the row label, 'Revenue' tells you nothing about whose revenue",
            "",
            "Challenge: Embedding models process TEXT. How do you represent a table as text?",
            "",
            "Bad approach (raw dump): '2400000 1800000 1200000' → no context, terrible embeddings",
            "Good approach (natural language): 'North America Q4 Revenue: $2.4M, Growth: +29%'",
            "",
            "Key insight: Each row should include its column headers to be self-contained",
            "The embedding for 'Q4 Revenue: $2.4M' is far more useful than '$2.4M' alone",
        ],
    )

    # SLIDE 37: Converting Tables
    add_content_slide(
        prs,
        "Converting Tables to RAG-Ready Text",
        [
            "Strategy 1: Natural Language Descriptions (RECOMMENDED)",
            "  'Employee Alice Johnson works in Engineering as Senior Developer earning $125K'",
            "",
            "Strategy 2: Key-Value Format",
            "  'Name: Alice Johnson; Department: Engineering; Title: Senior Developer; Salary: $125K'",
            "",
            "Strategy 3: Row-Based Chunks (with headers repeated)",
            "  [Employee Directory - rows 1-5] Columns: Name, Dept, Title",
            "  - Name: Alice; Dept: Engineering; Title: Senior Developer",
            "",
            "Strategy 4: Summary + Detail",
            "  Summary chunk: 'Employee directory with 8 records. Salary range: $75K-$145K'",
            "  Detail chunks: One chunk per row or per 3-5 rows",
            "",
            "Always include: Column headers, sheet name, and summary statistics as metadata",
        ],
    )

    # SLIDE 38: Spreadsheet Methods
    add_table_slide(
        prs,
        "Spreadsheet Methods Comparison",
        ["Method", "File Types", "Dependencies", "Type Inference", "Best For"],
        [
            [
                "openpyxl",
                "XLSX only",
                "Lightweight",
                "Basic",
                "Cell-level access, formatting, merged cells",
            ],
            [
                "pandas",
                "XLSX, CSV, XLS+",
                "Heavy (numpy)",
                "Excellent",
                "Data analysis, aggregation, stats",
            ],
            [
                "csv (stdlib)",
                "CSV only",
                "None",
                "None (strings)",
                "Zero dependencies, streaming large files",
            ],
        ],
        intro="Recommendation: pandas for analysis + extraction → openpyxl for formatting → csv for simple files",
    )

    # =================================================================
    # SECTION 8: Images & OCR
    # =================================================================
    add_section_divider(prs, "Section 8", "Images & OCR Deep Dive — From Pixels to Text")

    # SLIDE 39: What is an Image?
    add_content_slide(
        prs,
        "What is an Image? — Understanding Pixels",
        [
            "An image is a 2D grid (matrix) of tiny colored squares called PIXELS",
            "",
            "Each pixel has 3 color channels: Red, Green, Blue (RGB) — each from 0 to 255",
            "  (255, 0, 0) = pure red    (0, 255, 0) = pure green    (0, 0, 255) = pure blue",
            "  (0, 0, 0) = black         (255, 255, 255) = white      (128, 128, 128) = gray",
            "",
            "Image sizes: Width x Height in pixels",
            "  800 x 600 image = 480,000 pixels = 1,440,000 color values to process",
            "  A full page at 300 DPI ≈ 2,550 x 3,300 pixels = 8.4 million pixels",
            "",
            "Resolution (DPI = Dots Per Inch):",
            "  72 DPI = screen quality (blurry text for OCR)",
            "  150 DPI = reasonable quality",
            "  300 DPI = print quality (recommended for OCR)",
            "",
            "The key point: An image contains ZERO text data. Only pixel colors.",
        ],
    )

    # SLIDE 40: What is OCR?
    add_content_slide(
        prs,
        "What is OCR? — Optical Character Recognition",
        [
            "OCR is the process of converting images of text into actual machine-readable text",
            "",
            "The fundamental challenge:",
            "  Input: A grid of pixel colors → [0,0,0, 255,255,255, 0,0,0, ...]",
            "  Output: Characters → 'Hello World'",
            "",
            "How does the computer know that a group of dark pixels forms the letter 'A'?",
            "  Traditional OCR: Pattern matching against known character templates",
            "  Modern OCR: Deep learning models trained on millions of text images",
            "",
            "OCR must handle:",
            "  - Different fonts (serif, sans-serif, handwriting, monospace)",
            "  - Different sizes (6pt footnote to 72pt headline)",
            "  - Different qualities (clean print, faded photocopy, smartphone photo)",
            "  - Different languages (Latin, CJK, Arabic, Devanagari, etc.)",
            "  - Noise (smudges, creases, uneven lighting, background patterns)",
        ],
    )

    # SLIDE 41: OCR Pipeline
    add_content_slide(
        prs,
        "The OCR Pipeline: Step by Step",
        [
            "Step 1: LOAD IMAGE",
            "  Read the image file (PNG, JPEG, TIFF) into memory as a pixel array",
            "",
            "Step 2: PREPROCESS (Critical for accuracy!)",
            "  Convert to grayscale → Apply threshold (binarize) → Denoise → Deskew (straighten)",
            "",
            "Step 3: LAYOUT ANALYSIS",
            "  Detect text regions, columns, paragraphs, and line boundaries",
            "  Separate text from images, tables, and whitespace",
            "",
            "Step 4: CHARACTER RECOGNITION",
            "  For each detected text region: identify individual characters using ML models",
            "  Tesseract uses LSTM neural networks trained on millions of text samples",
            "",
            "Step 5: POST-PROCESSING",
            "  Apply dictionary/language model corrections, fix common OCR errors",
            "  Output: Plain text with optional bounding box coordinates",
        ],
    )

    # SLIDE 42: Preprocessing Deep Dive
    add_content_slide(
        prs,
        "Image Preprocessing for OCR — The Accuracy Multiplier",
        [
            "Preprocessing can improve OCR accuracy from 70% to 95%+ on noisy documents",
            "",
            "1. GRAYSCALE CONVERSION",
            "   RGB (3 channels) → single brightness channel. Simplifies analysis.",
            "",
            "2. BINARIZATION (Thresholding)",
            "   Every pixel becomes either black (text) or white (background).",
            "   Methods: Simple threshold, Otsu's method (auto), adaptive threshold (local)",
            "",
            "3. NOISE REMOVAL",
            "   Median filter removes salt-and-pepper noise from scans",
            "   Morphological operations (erode/dilate) clean up character edges",
            "",
            "4. DESKEWING",
            "   Detect and correct rotation from tilted scans (even 1-2 degrees matters!)",
            "",
            "5. SHARPENING",
            "   Enhance blurry text edges using unsharp mask or convolution filters",
        ],
    )

    # SLIDE 43: Tesseract OCR
    add_content_slide(
        prs,
        "Tesseract OCR — The Industry Standard",
        [
            "Tesseract: Open-source OCR engine by Google. 100+ language support.",
            "",
            "Key configuration — Page Segmentation Modes (PSM):",
            "  PSM 3 (default): Fully automatic page segmentation — best for full pages",
            "  PSM 6: Assume a single block of text — good for cropped regions",
            "  PSM 7: Treat image as a single text line",
            "  PSM 13: Raw line — for individual characters or special cases",
            "",
            "OEM (OCR Engine Modes):",
            "  OEM 1: LSTM neural network engine (default, most accurate)",
            "  OEM 0: Legacy Tesseract engine (faster but less accurate)",
            "",
            "Output formats:",
            "  String: Plain text output",
            "  HOCR: HTML-like output with bounding boxes for every word",
            "  Data: Word-by-word output with coordinates and confidence scores",
            "",
            "Typical accuracy: ~95% on clean 300 DPI scans, ~80% on noisy low-res images",
        ],
    )

    # SLIDE 44: EasyOCR
    add_content_slide(
        prs,
        "EasyOCR — Deep Learning Alternative",
        [
            "EasyOCR: PyTorch-based OCR with 80+ language support",
            "",
            "How it differs from Tesseract:",
            "  - Uses deep learning for both detection AND recognition",
            "  - CRAFT model for text detection (where is the text?)",
            "  - CRNN model for text recognition (what does it say?)",
            "  - Better with natural scene text (photos, signs, not just documents)",
            "",
            "Key features:",
            "  - Returns bounding boxes + text + confidence score for each detected region",
            "  - GPU acceleration with CUDA (much faster than CPU)",
            "  - Multiple language detection in same image",
            "  - No system dependencies (unlike Tesseract which needs separate install)",
            "",
            "Trade-offs:",
            "  - Larger install size (~1-2 GB with PyTorch)",
            "  - First run downloads model files (~100 MB per language)",
            "  - Slower than Tesseract on CPU for clean document scans",
        ],
    )

    # SLIDE 45: OCR Comparison
    add_table_slide(
        prs,
        "OCR Method Comparison — When to Use Each",
        [
            "Aspect",
            "Tesseract",
            "EasyOCR",
            "Cloud APIs (Google/AWS)",
            "Vision LLMs (GPT-4V)",
        ],
        [
            [
                "Type",
                "Traditional + LSTM",
                "Deep learning (PyTorch)",
                "Cloud-hosted DL",
                "Multimodal LLM",
            ],
            ["Install Size", "~30 MB", "~1-2 GB", "None (API)", "None (API)"],
            [
                "Speed (per page)",
                "0.5-2 sec",
                "1-5 sec (CPU)",
                "1-3 sec (network)",
                "2-10 sec",
            ],
            ["Clean Doc Accuracy", "~95%", "~93%", "~98%", "~97%"],
            ["Noisy Doc Accuracy", "~80%", "~85%", "~95%", "~93%"],
            ["Cost", "Free", "Free", "$1-3 per 1K pages", "$0.01-0.03 per image"],
            [
                "Best For",
                "Documents, forms",
                "Scene text, photos",
                "Production at scale",
                "Complex layouts",
            ],
        ],
    )

    # =================================================================
    # SECTION 9: Email, Markdown, EPUB
    # =================================================================
    add_section_divider(prs, "Section 9", "Email, Markdown & EPUB")

    # SLIDE 46: Email
    add_content_slide(
        prs,
        "Email Parsing for RAG",
        [
            "Email format uses MIME (Multipurpose Internet Mail Extensions) standard",
            "",
            "Extraction targets:",
            "  Headers: From, To, CC, Subject, Date — valuable metadata for filtering",
            "  Body: Plain text and/or HTML versions (prefer plain text; strip HTML if only HTML)",
            "  Attachments: PDFs, images, documents — parse these with their respective methods",
            "",
            "RAG-ready format for each email:",
            "  'Email from {from} to {to} on {date}. Subject: {subject}.",
            "   Body: {body_text}'",
            "",
            "Chunking strategy: Usually one email = one chunk (emails are typically short)",
            "For long email threads: Split by message, include thread metadata",
            "",
            "Common pitfalls: RFC 2047 encoded headers, nested MIME parts, inline HTML images",
            "Python's built-in email module handles all of this — no external dependencies needed",
        ],
    )

    # SLIDE 47: Markdown/Text
    add_content_slide(
        prs,
        "Markdown & Text: Chunking IS the Challenge",
        [
            "Markdown and plain text are the easiest to 'parse' — they're already text!",
            "The real challenge is choosing the right CHUNKING STRATEGY",
            "",
            "For Markdown (with headings):",
            "  Use heading-aware chunking — each section under a heading becomes one chunk",
            "  Extract code blocks separately — they need different treatment for RAG",
            "  Parse the AST (Abstract Syntax Tree) with mistune for structured extraction",
            "",
            "For Plain Text (no headings):",
            "  Paragraph-based: Split on double newlines (\\n\\n) — respects topic boundaries",
            "  Sentence-based: Group 3-5 sentences per chunk — good semantic coherence",
            "  Recursive splitting: Try paragraphs first, then sentences, then characters",
            "",
            "Key insight: Plain text without headings is actually the HARDEST to chunk well",
            "because there are no structural markers to guide the splitting",
        ],
    )

    # SLIDE 48: Text Chunking Challenge
    add_content_slide(
        prs,
        "Why Plain Text Chunking Needs Careful Thought",
        [
            "Consider this text about World War II — where do you split it?",
            "",
            "Character chunking at 500 chars: Might split mid-sentence about a battle",
            '  "...the Allied forces landed at Norman" | "dy on June 6, 1944..."',
            "  → Broken context, incomplete information in both chunks",
            "",
            "Sentence chunking (3 sentences): Better, but related sentences may be split",
            "  Chunk 1: Setup of D-Day planning. Chunk 2: The actual landing details.",
            "  → A query about 'D-Day preparation' might miss the actual event",
            "",
            "Paragraph chunking: Best for topic-organized text",
            "  Each paragraph usually covers one sub-topic → natural chunk boundary",
            "  But some paragraphs are very long, others very short",
            "",
            "Solution: Recursive splitting with overlap — tries paragraphs first, falls back to sentences",
        ],
    )

    # SLIDE 49: EPUB
    add_content_slide(
        prs,
        "EPUB Parsing: Natural Chapter Boundaries",
        [
            "EPUB is excellent for RAG — it has built-in document structure!",
            "",
            "Extraction with ebooklib + BeautifulSoup:",
            "  1. Read EPUB with ebooklib → get chapter items in reading order",
            "  2. Parse each chapter's XHTML content with BeautifulSoup",
            "  3. Extract clean text, headings, lists, and tables from each chapter",
            "  4. Preserve chapter structure as metadata",
            "",
            "Chunking strategies for EPUB:",
            "  Chapter-per-chunk: One chunk per chapter — simple and effective for short chapters",
            "  Heading-aware: Convert to markdown, then split by headings — best for long chapters",
            "  Recursive split: For very long chapters, fall back to recursive character splitting",
            "",
            "Metadata to include: Book title, author, chapter title, chapter number",
            "This enables queries like 'What does Chapter 3 of [book] discuss?'",
        ],
    )

    # SLIDE 50: Minor Formats Decision
    add_table_slide(
        prs,
        "Email, Markdown, EPUB — Quick Decision Guide",
        ["Format", "Go-To Method", "Chunking Strategy", "Key Tip"],
        [
            [
                "Email (.eml)",
                "Python email (stdlib)",
                "One email = one chunk",
                "Include headers as metadata",
            ],
            [
                "Markdown (.md)",
                "mistune AST parsing",
                "Heading-aware chunking",
                "Extract code blocks separately",
            ],
            [
                "Plain text (.txt)",
                "Direct read + chunk",
                "Recursive splitting",
                "Paragraph splits > sentence splits",
            ],
            [
                "EPUB (.epub)",
                "ebooklib + BS4",
                "Chapter-per-chunk",
                "Use spine for correct reading order",
            ],
        ],
    )

    # =================================================================
    # SECTION 10: Chunking Strategies
    # =================================================================
    add_section_divider(prs, "Section 10", "Chunking Strategies — Where RAG Quality is Won or Lost")

    # SLIDE 51: What is Chunking
    add_content_slide(
        prs,
        "What is Chunking & Why It Matters for RAG",
        [
            "Chunking = splitting extracted text into smaller pieces for embedding and retrieval",
            "",
            "Why chunk? Three reasons:",
            "  1. Embedding models have token limits (e.g., 512 tokens for many models)",
            "  2. Smaller chunks have more focused semantic meaning → better retrieval",
            "  3. LLM context windows are limited — you want to inject only relevant content",
            "",
            "The chunk size paradox:",
            "  Too small: Each chunk lacks context. 'The revenue was $2.4M' — whose revenue?",
            "  Too large: Semantic meaning is diluted. A 5-page chunk matches too many queries.",
            "  Sweet spot: 200-500 tokens per chunk for most embedding models",
            "",
            "Chunking is arguably the MOST important step in the RAG pipeline",
            "Bad chunking → bad embeddings → bad retrieval → bad LLM answers",
            "The right chunking strategy depends on the document type and structure",
        ],
    )

    # SLIDE 52: Character Chunking
    add_content_slide(
        prs,
        "Strategy 1: Fixed Character Chunking",
        [
            "How it works: Slide a fixed-size window across the text with overlap",
            "",
            "Example (chunk_size=200, overlap=50):",
            '  Chunk 1: "Artificial Intelligence is a branch of computer science that..."  [chars 0-200]',
            '  Chunk 2: "...science that aims to create systems capable of performing..."  [chars 150-350]',
            '  Chunk 3: "...performing tasks that normally require human intelligence..."  [chars 300-500]',
            "",
            "Pros:",
            "  + Predictable, uniform chunk sizes → consistent token counts",
            "  + Very simple to implement",
            "  + Fast execution",
            "",
            "Cons:",
            "  - Splits mid-word and mid-sentence → broken context",
            "  - No awareness of document structure",
            '  - "...the Allied forces landed at Norman" | "dy on June 6, 1944..."',
            "",
            "Best for: Baseline comparison, when uniformity matters more than quality",
        ],
    )

    # SLIDE 53: Sentence Chunking
    add_content_slide(
        prs,
        "Strategy 2: Sentence-Based Chunking",
        [
            "How it works: Split text into sentences, then group N sentences per chunk",
            "",
            "Sentence detection: Split on '.', '!', '?' followed by whitespace",
            "  (Can fail on abbreviations: 'Dr. Smith' or numbers: '3.14')",
            "",
            "Example (sentences_per_chunk=3, overlap=1):",
            "  Chunk 1: [Sentence 1] [Sentence 2] [Sentence 3]",
            "  Chunk 2: [Sentence 3] [Sentence 4] [Sentence 5]   ← Sentence 3 appears in both!",
            "  Chunk 3: [Sentence 5] [Sentence 6] [Sentence 7]",
            "",
            "Pros:",
            "  + Preserves complete sentences → better semantic coherence",
            "  + Overlap at sentence boundaries → context preserved at splits",
            "",
            "Cons:",
            "  - Variable chunk sizes (short vs long sentences)",
            "  - No awareness of paragraphs or headings",
            "  - Sentence detection can fail on edge cases",
        ],
    )

    # SLIDE 54: Recursive Splitting
    add_content_slide(
        prs,
        "Strategy 3: Recursive Character Splitting",
        [
            "How it works: Try the most meaningful separator first, fall back to simpler ones",
            "",
            "Separator hierarchy:",
            "  1st try: Split on PARAGRAPHS (double newline \\n\\n)",
            "  2nd try: Split on NEWLINES (single \\n)",
            "  3rd try: Split on SENTENCES ('. ')",
            "  4th try: Split on SPACES (' ')",
            "  Last resort: Split on CHARACTERS (hard cut)",
            "",
            "At each level: merge parts until they approach chunk_size, then split",
            "",
            "This is the most popular strategy (used by LangChain's RecursiveCharacterTextSplitter)",
            "",
            "Pros:",
            "  + Respects document structure at multiple levels",
            "  + Graceful degradation — uses the best boundary available",
            "  + Good balance of size consistency and semantic coherence",
            "",
            "Cons:",
            "  - Still no heading awareness (treats all paragraphs equally)",
            "  - More complex implementation than simple strategies",
        ],
    )

    # SLIDE 55: Heading-Aware Chunking
    add_content_slide(
        prs,
        "Strategy 4: Heading-Aware Chunking (RECOMMENDED)",
        [
            "How it works: Split on headings, use heading text as metadata for each chunk",
            "",
            "Example input (markdown):",
            "  # Introduction     → Chunk 1: heading='Introduction', content='...'",
            "  paragraph text...",
            "  # Methods           → Chunk 2: heading='Methods', content='...'",
            "  paragraph text...",
            "  ## Experiment Setup  → Chunk 3: heading='Experiment Setup', content='...'",
            "",
            "Each chunk carries its heading as metadata → enables section-based retrieval",
            "",
            "Pros:",
            "  + Respects the author's intended topic structure",
            "  + Headings provide natural, meaningful boundaries",
            "  + Metadata enables filtered queries ('find info in Methods section')",
            "",
            "Cons:",
            "  - Requires documents with headings (doesn't work for plain text)",
            "  - Variable chunk sizes (some sections are very short, others very long)",
            "  - Very long sections may need secondary splitting",
        ],
    )

    # SLIDE 56: Chunk Size Impact
    add_content_slide(
        prs,
        "Chunk Size Impact on RAG Quality",
        [
            "The chunk size you choose directly affects retrieval quality:",
            "",
            "TOO SMALL (50-100 tokens):",
            "  \"The revenue was $2.4M\" → Matches 'revenue' queries but lacks context",
            "  Who? When? What division? The LLM has no context to generate a good answer.",
            "",
            "TOO LARGE (2000+ tokens):",
            "  Full page of mixed content → Matches many queries but most content is irrelevant",
            "  The embedding becomes an 'average' of many topics → diluted semantic meaning.",
            "",
            "SWEET SPOT (200-500 tokens):",
            '  "In Q4 2024, North America revenue reached $2.4M, representing 29% year-over-year',
            '   growth driven primarily by enterprise suite sales."',
            "  Complete thought with context, focused enough for precise retrieval.",
            "",
            "OVERLAP (10-20% of chunk size):",
            "  Repeating 50-100 tokens at chunk boundaries prevents information loss at splits.",
        ],
    )

    # SLIDE 57: Chunking per Document Type
    add_table_slide(
        prs,
        "Recommended Chunking Strategy per Document Type",
        ["Document Type", "Primary Strategy", "Fallback Strategy", "Chunk Size"],
        [
            [
                "PDF (with headings)",
                "Heading-aware (detect by font size)",
                "Recursive split",
                "300-500 tokens",
            ],
            ["PDF (plain text)", "Recursive split", "Sentence-based", "400-600 tokens"],
            [
                "DOCX",
                "Heading-aware (style names)",
                "Recursive split",
                "300-500 tokens",
            ],
            [
                "PPTX",
                "Slide-per-chunk",
                "Sentence-based on merged text",
                "1 slide = 1 chunk",
            ],
            [
                "HTML",
                "Heading-aware (via html2text)",
                "Recursive split",
                "300-500 tokens",
            ],
            [
                "Spreadsheets",
                "Row-based (N rows per chunk)",
                "Natural language per row",
                "5-10 rows",
            ],
            ["Images (OCR)", "Paragraph-based", "Sentence-based", "200-400 tokens"],
            ["Email", "One email per chunk", "Thread-based", "Full email"],
            ["Markdown", "Heading-aware", "Paragraph-based", "300-500 tokens"],
            ["EPUB", "Chapter-per-chunk", "Heading-aware within chapters", "1 chapter"],
        ],
        intro="Match your chunking strategy to the document's natural structure.",
    )

    # SLIDE 58: Chunking Best Practices
    add_content_slide(
        prs,
        "Chunking Best Practices — 8 Key Rules",
        [
            "1. Always include overlap (10-20% of chunk size) to preserve context at boundaries",
            "2. Include metadata with every chunk: source file, page/section, heading, date",
            "3. Test with real queries — the best strategy depends on how users will search",
            "4. Don't mix document types in the same chunk — keep chunks from one source",
            "5. Extract tables separately — convert to natural language, don't dump raw cells",
            "6. For structured docs (DOCX, HTML, MD): heading-aware chunking almost always wins",
            "7. For unstructured docs (plain text, OCR output): recursive splitting is the best default",
            "8. Monitor chunk size distribution — very skewed distributions indicate a poor strategy",
        ],
    )

    # =================================================================
    # SECTION 11: Putting It All Together
    # =================================================================
    add_section_divider(prs, "Section 11", "Putting It All Together — Decision Matrix & Best Practices")

    # SLIDE 59: Complete Decision Matrix
    add_table_slide(
        prs,
        "Complete Decision Matrix: Document Type → Parser → Strategy",
        ["Document", "Recommended Parser", "Chunking Strategy", "Key Consideration"],
        [
            [
                "PDF (text)",
                "PyMuPDF or pypdf",
                "Recursive split",
                "Try direct extraction first",
            ],
            [
                "PDF (tables)",
                "pdfplumber",
                "Table-per-chunk + NL conversion",
                "Always include column headers",
            ],
            [
                "PDF (scanned)",
                "pytesseract",
                "Paragraph/sentence split",
                "Preprocess images for accuracy",
            ],
            [
                "DOCX",
                "python-docx",
                "Heading-aware",
                "Use style names for heading detection",
            ],
            [
                "PPTX",
                "python-pptx (structured)",
                "Slide-per-chunk",
                "Always include speaker notes",
            ],
            [
                "HTML (article)",
                "Trafilatura",
                "Heading-aware",
                "Remove boilerplate first",
            ],
            [
                "HTML (known site)",
                "BeautifulSoup",
                "Section-based",
                "Target <article> or <main>",
            ],
            [
                "XLSX",
                "pandas or openpyxl",
                "Row-based + NL",
                "Include headers with every row",
            ],
            ["CSV", "csv stdlib or pandas", "Row-based + NL", "Detect encoding first"],
            [
                "Images",
                "Tesseract (+preprocess)",
                "Paragraph split",
                "Always preprocess for quality",
            ],
            ["Email", "email stdlib", "One per chunk", "Include subject + metadata"],
            [
                "Markdown",
                "mistune AST",
                "Heading-aware",
                "Extract code blocks separately",
            ],
            [
                "EPUB",
                "ebooklib + BS4",
                "Chapter-per-chunk",
                "Follow spine reading order",
            ],
        ],
    )

    # SLIDE 60: Pipeline in Practice
    add_content_slide(
        prs,
        "The RAG Parsing Pipeline in Practice — End-to-End Example",
        [
            "Scenario: Parse a company's quarterly report (PDF with tables and charts)",
            "",
            "Step 1: Detect PDF type → Try pypdf extraction → Text is present → Born-digital PDF",
            "Step 2: Extract text with PyMuPDF (fast) → Get 15 pages of text",
            "Step 3: Extract tables with pdfplumber → 4 tables with financial data",
            "Step 4: Convert tables to natural language: 'Q4 Revenue for North America: $3.1M (+29% YoY)'",
            "Step 5: Detect headings by font size analysis → 'Executive Summary', 'Financial Results', etc.",
            "Step 6: Heading-aware chunking → 12 text chunks + 4 table chunks",
            "Step 7: Add metadata to each chunk: {source: 'Q4_Report.pdf', section: '...', page: N}",
            "Step 8: Embed each chunk → 16 vectors stored in vector database",
            "",
            "Query: 'What was the Q4 revenue growth in North America?'",
            "Retrieved chunk: 'Q4 Revenue for North America: $3.1M, representing +29% YoY growth...'",
            "LLM generates: 'North America Q4 revenue grew 29% year-over-year to $3.1M.'",
        ],
    )

    # SLIDE 61: Top 10 Pitfalls
    add_content_slide(
        prs,
        "Top 10 Common Pitfalls Across All Document Types",
        [
            "1. Not removing boilerplate from HTML → 'Home | About | Contact' in every chunk",
            "2. Dumping table cells as raw text → '$2.4M' without headers is meaningless",
            "3. Using OCR on born-digital PDFs → 100x slower with worse accuracy",
            "4. Ignoring speaker notes in PPTX → missing the best text for RAG",
            "5. Splitting mid-sentence with character chunking → broken context in embeddings",
            "6. No overlap between chunks → important context lost at boundaries",
            "7. Mixing boilerplate (headers/footers) with content in PDF → noise in every chunk",
            "8. Not including column headers with spreadsheet rows → context-free numbers",
            "9. Using one method for all PDFs → tables need pdfplumber, text needs PyMuPDF",
            "10. Skipping image preprocessing before OCR → 80% accuracy vs 95% with preprocessing",
        ],
    )

    # SLIDE 62: Production Tips
    add_content_slide(
        prs,
        "Production Tips for Document Parsing at Scale",
        [
            "1. Auto-detect document type by extension AND content analysis (don't trust extensions alone)",
            "2. Build a fallback chain: Try method A → if poor results → try method B → try OCR",
            "3. Monitor extraction quality: Track empty chunks, very short chunks, encoding errors",
            "4. Cache parsed results — re-parsing is expensive, especially OCR",
            "5. Process documents in parallel — most parsing libraries are single-threaded",
            "6. Set timeouts for OCR — a corrupted image can hang Tesseract indefinitely",
            "7. Store raw extracted text alongside chunks — makes re-chunking possible without re-parsing",
            "8. Version your chunks — when you change chunking strategy, you need to re-embed",
            "9. Log which method was used for each document — debugging requires knowing the extraction path",
            "10. Test with representative documents from your actual corpus, not just clean samples",
        ],
    )

    # SLIDE 63: Recommended Approach Summary
    add_table_slide(
        prs,
        "Recommended Approach Summary — Per Document Type",
        ["Document Type", "First Choice", "Second Choice", "Chunking"],
        [
            [
                "PDF (general)",
                "PyMuPDF (speed)",
                "pdfplumber (accuracy)",
                "Heading-aware / recursive",
            ],
            [
                "PDF (tables)",
                "pdfplumber",
                "PyMuPDF + custom",
                "Table-to-NL + heading-aware",
            ],
            [
                "PDF (scanned)",
                "Tesseract + preprocess",
                "EasyOCR / Cloud API",
                "Paragraph / sentence",
            ],
            ["Word (DOCX)", "python-docx", "mammoth (markdown)", "Heading-aware"],
            [
                "PowerPoint",
                "python-pptx (structured)",
                "python-pptx (basic)",
                "Slide-per-chunk",
            ],
            [
                "HTML / Web",
                "Trafilatura (articles)",
                "BeautifulSoup (custom)",
                "Heading-aware",
            ],
            [
                "Spreadsheets",
                "pandas (analysis)",
                "openpyxl (formatting)",
                "Row-based NL conversion",
            ],
            [
                "Images/Scans",
                "Tesseract (clean docs)",
                "EasyOCR (scene text)",
                "Paragraph / sentence",
            ],
            ["Email", "email stdlib", "—", "One email = one chunk"],
            ["Markdown", "mistune AST", "Simple regex", "Heading-aware"],
            ["Plain Text", "Direct read", "—", "Recursive splitting"],
            ["EPUB", "ebooklib + BS4", "—", "Chapter-per-chunk"],
        ],
    )

    # SLIDE 64: Key Takeaways
    add_content_slide(
        prs,
        "Key Takeaways",
        [
            "1. Every document format stores data differently internally — one parser doesn't fit all",
            "",
            "2. PDF is the most complex: 6 methods exist because no single one handles everything",
            "",
            "3. Tables need special treatment: Always convert to natural language with headers",
            "",
            "4. OCR is a last resort: Only use it for image-based documents, always preprocess first",
            "",
            (
                "5. Chunking strategy matters more than most people think:"
                " heading-aware > recursive > sentence > character"
            ),
            "",
            "6. Metadata is as important as content: Source, section, page number enable filtered retrieval",
            "",
            "7. Test with real queries from real users: The best strategy depends on search patterns",
            "",
            "8. Quality of parsing = Quality of RAG: Garbage in, garbage out applies strongly here",
        ],
    )

    # SLIDE 65: Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "Questions & Discussion"
    p2.font.size = Pt(28)
    p2.font.color.rgb = LIGHT_BLUE
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(30)

    p3 = tf.add_paragraph()
    p3.text = "All code, samples, and guides available in the repository"
    p3.font.size = Pt(18)
    p3.font.color.rgb = RGBColor(0x90, 0xB0, 0xD0)
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(40)

    return prs


# =========================================================================
# Main
# =========================================================================

if __name__ == "__main__":
    print("Generating presentation...")
    prs = build_presentation()
    prs.save(str(OUTPUT_PATH))
    print(f"\nCreated: {OUTPUT_PATH}")
    print(f"Total slides: {len(prs.slides)}")
    print("\nOpen the .pptx file in PowerPoint or Google Slides to view.")
